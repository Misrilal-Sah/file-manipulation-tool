"""
Document Conversion Service
Supports PDF <-> Word conversions using LibreOffice and PyMuPDF
"""
import subprocess
import os
import tempfile
from pathlib import Path
from typing import Optional, List, Tuple
import fitz  # PyMuPDF
from PIL import Image
import io

from app.config import settings
from app.utils.helpers import generate_file_id


class ConvertService:
    """Service for document format conversion"""
    
    @staticmethod
    def check_libreoffice() -> bool:
        """Check if LibreOffice is available"""
        return Path(settings.LIBREOFFICE_PATH).exists()
    
    @staticmethod
    def word_to_pdf_libreoffice(file_path: Path, original_filename: str = None) -> Tuple[Path, bool]:
        """
        Convert Word document to PDF using LibreOffice
        Returns: (output_path, success)
        """
        if not ConvertService.check_libreoffice():
            return file_path, False
        
        output_dir = settings.OUTPUT_DIR
        
        # Generate meaningful filename
        if original_filename:
            base_name = Path(original_filename).stem
            output_filename = f"{base_name}.pdf"
        else:
            output_filename = f"{generate_file_id()}_converted.pdf"
        
        try:
            # Run LibreOffice headless conversion
            cmd = [
                settings.LIBREOFFICE_PATH,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', str(output_dir),
                str(file_path)
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            
            if result.returncode == 0:
                # Find the converted file (LibreOffice uses original name)
                original_name = file_path.stem
                converted_file = output_dir / f"{original_name}.pdf"
                
                if converted_file.exists():
                    # Rename to our desired filename
                    final_path = output_dir / output_filename
                    # If file exists with same name, add unique suffix
                    if final_path.exists() and final_path != converted_file:
                        final_path = output_dir / f"{Path(output_filename).stem}_{generate_file_id()[:8]}.pdf"
                    converted_file.rename(final_path)
                    return final_path, True
            
            return file_path, False
        except Exception as e:
            print(f"Conversion error: {e}")
            return file_path, False
    
    @staticmethod
    def pdf_to_word_basic(file_path: Path, original_filename: str = None) -> Path:
        """
        Convert PDF to Word (basic text extraction)
        Note: This preserves text but not complex layouts
        """
        from docx import Document
        
        # Generate meaningful filename
        if original_filename:
            base_name = Path(original_filename).stem
            output_filename = f"{base_name}.docx"
        else:
            output_filename = f"{generate_file_id()}_converted.docx"
        output_path = settings.OUTPUT_DIR / output_filename
        
        # Open PDF and extract text
        pdf_doc = fitz.open(file_path)
        word_doc = Document()
        
        try:
            for page_num, page in enumerate(pdf_doc):
                # Extract text blocks
                text = page.get_text()
                
                if page_num > 0:
                    word_doc.add_page_break()
                
                # Add text to Word document
                for line in text.split('\n'):
                    if line.strip():
                        word_doc.add_paragraph(line)
            
            word_doc.save(output_path)
        finally:
            pdf_doc.close()
        
        return output_path
    
    @staticmethod
    def pdf_to_images(file_path: Path, dpi: int = 200, 
                      image_format: str = "png", original_filename: str = None) -> List[tuple]:
        """
        Convert PDF pages to images
        Returns: List of tuples (storage_path, display_filename)
        """
        pdf_doc = fitz.open(file_path)
        output_results = []
        
        # Convert format for PIL (jpg -> JPEG)
        pil_format = "JPEG" if image_format.lower() in ["jpg", "jpeg"] else image_format.upper()
        file_extension = "jpg" if image_format.lower() in ["jpg", "jpeg"] else image_format.lower()
        
        # Get base name for display filenames
        if original_filename:
            base_name = Path(original_filename).stem
        else:
            base_name = "page"
        
        # Generate a unique batch ID to prevent collisions in storage
        batch_id = generate_file_id()[:8]
        
        try:
            total_pages = len(pdf_doc)
            for page_num, page in enumerate(pdf_doc):
                mat = fitz.Matrix(dpi / 72, dpi / 72)
                pix = page.get_pixmap(matrix=mat)
                
                # Convert to PIL Image
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                
                # Storage filename (with batch_id for uniqueness)
                storage_filename = f"{batch_id}_{base_name}_{page_num + 1}.{file_extension}"
                output_path = settings.OUTPUT_DIR / storage_filename
                
                # Display filename (clean, user-friendly)
                if total_pages == 1:
                    display_filename = f"{base_name}.{file_extension}"
                else:
                    display_filename = f"{base_name}({page_num + 1}).{file_extension}"
                
                img.save(output_path, format=pil_format)
                output_results.append((output_path, display_filename))
        finally:
            pdf_doc.close()
        
        return output_results
    
    @staticmethod
    def images_to_pdf(image_paths: List[Path], output_filename: str = "combined.pdf") -> Path:
        """
        Combine multiple images into a PDF
        """
        output_path = settings.OUTPUT_DIR / f"{generate_file_id()}_{output_filename}"
        
        pdf_doc = fitz.open()
        
        try:
            for image_path in image_paths:
                # Open image and get dimensions
                img = Image.open(image_path)
                if img.mode != "RGB":
                    img = img.convert("RGB")
                
                # Create PDF page with image dimensions
                width, height = img.size
                
                # Convert to bytes
                img_bytes = io.BytesIO()
                img.save(img_bytes, format="PNG")
                img_bytes.seek(0)
                
                # Create page and insert image
                page = pdf_doc.new_page(width=width, height=height)
                page.insert_image(page.rect, stream=img_bytes.getvalue())
            
            pdf_doc.save(output_path)
        finally:
            pdf_doc.close()
        
        return output_path
    
    @staticmethod
    def excel_to_pdf(file_path: Path, original_filename: str = None) -> Tuple[Path, bool]:
        """
        Convert Excel to PDF using LibreOffice
        """
        return ConvertService._libreoffice_convert(file_path, "pdf", original_filename)
    
    @staticmethod
    def powerpoint_to_pdf(file_path: Path, original_filename: str = None) -> Tuple[Path, bool]:
        """
        Convert PowerPoint to PDF using LibreOffice
        """
        return ConvertService._libreoffice_convert(file_path, "pdf", original_filename)
    
    @staticmethod
    def html_to_pdf(file_path: Path, original_filename: str = None) -> Tuple[Path, bool]:
        """
        Convert HTML to PDF using LibreOffice
        """
        return ConvertService._libreoffice_convert(file_path, "pdf", original_filename)
    
    @staticmethod
    def _libreoffice_convert(file_path: Path, target_format: str, original_filename: str = None) -> Tuple[Path, bool]:
        """
        Generic LibreOffice conversion
        """
        if not ConvertService.check_libreoffice():
            return file_path, False
        
        output_dir = settings.OUTPUT_DIR
        
        try:
            cmd = [
                settings.LIBREOFFICE_PATH,
                '--headless',
                '--convert-to', target_format,
                '--outdir', str(output_dir),
                str(file_path)
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            
            if result.returncode == 0:
                # LibreOffice uses the stored filename (UUID-based) for output
                stored_name = file_path.stem
                converted_file = output_dir / f"{stored_name}.{target_format}"
                
                if converted_file.exists():
                    # Determine final filename
                    if original_filename:
                        # Use the user's original filename (e.g., "Report.pptx" -> "Report.pdf")
                        base_name = Path(original_filename).stem
                        final_filename = f"{base_name}.{target_format}"
                    else:
                        # Fallback: use stored name with unique ID
                        final_filename = f"{generate_file_id()[:8]}_{stored_name}.{target_format}"
                    
                    final_path = output_dir / final_filename
                    
                    # Handle potential name collision
                    if final_path.exists() and final_path != converted_file:
                        base_name = Path(final_filename).stem
                        final_path = output_dir / f"{base_name}_{generate_file_id()[:8]}.{target_format}"
                    
                    converted_file.rename(final_path)
                    return final_path, True
            
            return file_path, False
        except Exception as e:
            print(f"LibreOffice conversion error: {e}")
            return file_path, False
    
    @staticmethod
    def get_supported_conversions() -> dict:
        """
        Get list of supported conversions
        """
        libreoffice_available = ConvertService.check_libreoffice()
        
        return {
            "pdf_to_word": True,  # Basic text extraction always available
            "pdf_to_images": True,
            "images_to_pdf": True,
            "word_to_pdf": libreoffice_available,
            "excel_to_pdf": libreoffice_available,
            "powerpoint_to_pdf": libreoffice_available,
            "libreoffice_available": libreoffice_available
        }
    
    @staticmethod
    def pdf_to_powerpoint(file_path: Path, original_filename: str = None) -> Tuple[Path, str]:
        """
        Convert PDF to PowerPoint - each page becomes a slide with embedded image
        Returns: (storage_path, display_filename)
        """
        from pptx import Presentation
        from pptx.util import Inches
        
        pdf_doc = fitz.open(file_path)
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Get display filename
        if original_filename:
            base_name = Path(original_filename).stem
        else:
            base_name = "presentation"
        
        batch_id = generate_file_id()[:8]
        storage_filename = f"{batch_id}_{base_name}.pptx"
        display_filename = f"{base_name}.pptx"
        output_path = settings.OUTPUT_DIR / storage_filename
        
        try:
            for page_num, page in enumerate(pdf_doc):
                # Convert page to image
                mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for better quality
                pix = page.get_pixmap(matrix=mat)
                
                # Save temp image
                temp_img_path = settings.OUTPUT_DIR / f"temp_{batch_id}_{page_num}.png"
                pix.save(str(temp_img_path))
                
                # Add slide with image
                blank_layout = prs.slide_layouts[6]  # blank layout
                slide = prs.slides.add_slide(blank_layout)
                
                # Add image to slide (centered, scaled to fit)
                slide.shapes.add_picture(str(temp_img_path), Inches(0.5), Inches(0.5), 
                                         width=Inches(9), height=Inches(6.5))
                
                # Clean up temp image
                temp_img_path.unlink()
            
            prs.save(str(output_path))
        finally:
            pdf_doc.close()
        
        return output_path, display_filename
    
    @staticmethod
    def pdf_to_excel(file_path: Path, original_filename: str = None) -> Tuple[Path, str]:
        """
        Convert PDF to Excel - extracts text from each page into separate sheets
        Returns: (storage_path, display_filename)
        """
        import openpyxl
        
        pdf_doc = fitz.open(file_path)
        wb = openpyxl.Workbook()
        
        # Remove default sheet
        if "Sheet" in wb.sheetnames:
            del wb["Sheet"]
        
        # Get display filename
        if original_filename:
            base_name = Path(original_filename).stem
        else:
            base_name = "spreadsheet"
        
        batch_id = generate_file_id()[:8]
        storage_filename = f"{batch_id}_{base_name}.xlsx"
        display_filename = f"{base_name}.xlsx"
        output_path = settings.OUTPUT_DIR / storage_filename
        
        try:
            for page_num, page in enumerate(pdf_doc):
                # Create sheet for each page
                sheet_name = f"Page {page_num + 1}"[:31]  # Excel max sheet name is 31 chars
                ws = wb.create_sheet(sheet_name)
                
                # Get text from page
                text = page.get_text()
                
                # Split text into lines and put each line in a row
                lines = text.split('\n')
                for row_num, line in enumerate(lines, 1):
                    if line.strip():  # Skip empty lines
                        ws.cell(row=row_num, column=1, value=line)
            
            # Ensure at least one sheet exists
            if not wb.sheetnames:
                ws = wb.create_sheet("Page 1")
                ws.cell(row=1, column=1, value="(No text content found)")
            
            wb.save(str(output_path))
        finally:
            pdf_doc.close()
        
        return output_path, display_filename


# Create singleton instance
convert_service = ConvertService()

